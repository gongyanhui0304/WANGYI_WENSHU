"""Generate PPT pages from ppt/pages contracts.

Only each page folder is used:
- page.md: business rules
- config.json: program-readable page config
- notes.md: page notes
- example.png: the only visual/template reference for that page

Every enabled page with a valid example.png is included. No other visual source is used.
Before writing PPT, trace artifacts must be generated successfully.
"""
from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from validate_ppt_pages import validate
from generate_trace_artifacts import build_trace_artifacts

BASE = Path(__file__).resolve().parent.parent
PAGES = BASE / "ppt" / "pages"
OUT = BASE / "output" / "信息中心周报_页面契约版.pptx"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def load_config(page_dir: Path) -> dict:
    return json.loads((page_dir / "config.json").read_text(encoding="utf-8-sig"))


def configured_page_dirs() -> list[Path]:
    selected: list[Path] = []
    for page_dir in sorted(p for p in PAGES.iterdir() if p.is_dir()):
        config = load_config(page_dir)
        example = page_dir / "example.png"
        if config.get("enabled") is True and example.exists() and example.stat().st_size > 0:
            selected.append(page_dir)
    if not selected:
        raise RuntimeError("没有可生成页面：请检查 ppt/pages/<page>/config.json 的 enabled 和 example.png")
    return selected


def add_example_slide(prs: Presentation, page_dir: Path) -> None:
    config = load_config(page_dir)
    example = page_dir / "example.png"
    page_md = page_dir / "page.md"
    notes_md = page_dir / "notes.md"

    if not example.exists() or example.stat().st_size <= 0:
        raise RuntimeError(f"缺少有效样例图: {example}")
    if not page_md.exists() or not notes_md.exists():
        raise RuntimeError(f"页面契约不完整: {page_dir}")
    if config.get("visual_reference") != "example.png":
        raise RuntimeError(f"页面视觉参考必须是 example.png: {page_dir / 'config.json'}")
    if config.get("standard_ref") != "ppt/spec/page_standard.md":
        raise RuntimeError(f"页面未引用统一规范: {page_dir / 'config.json'}")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(example), 0, 0, width=SLIDE_W, height=SLIDE_H)


def main() -> None:
    problems = validate()
    if problems:
        raise SystemExit("页面契约校验失败，禁止生成：\n" + "\n".join(problems))

    try:
        run_dir = build_trace_artifacts(strict=True)
    except Exception as exc:
        if OUT.exists():
            OUT.unlink()
        raise SystemExit(str(exc))

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    pages = configured_page_dirs()
    for page_dir in pages:
        add_example_slide(prs, page_dir)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    names = ", ".join(page.name for page in pages)
    print(f"Trace artifacts: {run_dir}")
    print(f"Generated {OUT} ({len(prs.slides)} slides): {names}")

if __name__ == "__main__":
    main()

